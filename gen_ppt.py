"""Generate defence PPT from SLIDES.md content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette (matched to project CSS variables) ──
PHOSPHOR = RGBColor(0xE6, 0x7E, 0x22)  # orange accent
COBALT = RGBColor(0x08, 0x91, 0xB2)  # blue accent
INK = RGBColor(0x1A, 0x1A, 0x2E)  # dark text
INK_DIM = RGBColor(0x55, 0x55, 0x77)  # dim text
BG_PANEL = RGBColor(0xFA, 0xFB, 0xFD)  # light bg
BORDER = RGBColor(0xE2, 0xE4, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xE6, 0x7E, 0x22)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
GRAY_LIGHT = RGBColor(0xF0, 0xF1, 0xF3)

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

# ── Helper functions ──


def add_bg(slide, color=BG_PANEL):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left=0, top=0, width=None, height=Pt(4), color=PHOSPHOR):
    """Top accent line."""
    if width is None:
        width = prs.slide_width
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt_box(slide, left, top, width, height, text, font_size=Pt(18),
            color=INK, bold=False, alignment=PP_ALIGN.LEFT, font_name='宋体'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def multi_txt(slide, left, top, width, height, lines, font_size=Pt(14),
              color=INK, line_spacing=Pt(22), font_name='宋体'):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, sz, clr, bld) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = sz or font_size
        p.font.color.rgb = clr or color
        p.font.bold = bld or False
        p.font.name = font_name
        p.space_after = line_spacing
    return tf


def add_page_number(slide, num):
    txt_box(slide, 12.5, 7.0, 0.7, 0.4, str(num), Pt(10), INK_DIM, alignment=PP_ALIGN.RIGHT)


def add_section_title(slide, title, subtitle=None):
    add_accent_bar(slide)
    txt_box(slide, 0.8, 0.4, 11.5, 0.8, title, Pt(32), INK, bold=True, font_name='微软雅黑')
    if subtitle:
        txt_box(slide, 0.8, 1.1, 11.5, 0.5, subtitle, Pt(14), INK_DIM, font_name='微软雅黑')


def add_table(slide, left, top, col_widths, headers, rows, font_size=Pt(12)):
    """Add a styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                         Inches(left), Inches(top),
                                         Inches(total_w), Inches(0.4 * n_rows))
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)

    # Header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = hdr
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = '微软雅黑'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.color.rgb = INK
                p.font.name = '宋体'
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else GRAY_LIGHT

    return table_shape


def add_code_block(slide, left, top, width, height, code_text, font_size=Pt(11)):
    """Add a code block with dark background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(8)
    for i, line_text in enumerate(code_text.strip().split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = font_size
        p.font.color.rgb = RGBColor(0xCD, 0xD6, 0xF4)  # catppuccin text
        p.font.name = '宋体'
        p.space_after = Pt(2)
    return shape


def add_bullet_card(slide, left, top, width, height, title, bullets, title_color=PHOSPHOR):
    """Card with title + bullet points."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(10)
    tf.margin_right = Pt(10)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = '微软雅黑'
    p.space_after = Pt(8)

    for b in bullets:
        p = tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(12)
        p.font.color.rgb = INK
        p.font.name = '宋体'
        p.space_after = Pt(6)


def add_flow_arrow(slide, left, top, width, height, text):
    """Horizontal arrow shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PHOSPHOR
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER
    return shape


# ══════════════════════════════════════════
# P1 · Cover
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, INK)
add_accent_bar(slide, top=Inches(3.0), height=Pt(6), color=PHOSPHOR)
add_accent_bar(slide, top=Inches(3.15), height=Pt(2), color=COBALT)

txt_box(slide, 1.5, 1.2, 10.3, 1.0, 'OS I/O 软件层', Pt(48), WHITE, bold=True, font_name='微软雅黑')
txt_box(slide, 1.5, 1.9, 10.3, 0.8, 'read 操作模拟', Pt(40), PHOSPHOR, bold=True, font_name='微软雅黑')
txt_box(slide, 1.5, 3.8, 10.3, 0.6, '选题七 · 操作系统课程设计', Pt(18), INK_DIM)
txt_box(slide, 1.5, 4.4, 10.3, 0.5, 'I/O 软件层 sys_read() 全链路可视化教学工具  ·  8人小组', Pt(14), RGBColor(0x88, 0x88, 0xAA))
txt_box(slide, 1.5, 5.5, 10.3, 0.5, 'https://github.com/smithluo11/os-read-', Pt(12), COBALT)

# ══════════════════════════════════════════
# P2 · 课设要求
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '课设要求', '选题七：I/O 软件层 read 操作模拟 — 实验指导书要求')
add_page_number(slide, 2)

headers = ['序号', '功能要求', '我们的实现']
col_w = [0.6, 4.5, 6.5]
rows = [
    ['①', 'I/O 软件层分层模拟\n（用户层 → 设备无关 → 驱动 → 中断）',
     'L0-L4 五层拓扑图 + SVG 连线 + IRP↓ / IRQ↑ / DATA↑ 标签'],
    ['②', 'read 全流程：请求转换 → 权限校验\n→ 指令下发 → 中断响应 → 结果返回',
     '状态机引擎逐步推进，每步附 stepDescription，进程状态 RUNNING→BLOCKED→RUNNING'],
    ['③', '数据流转与缓冲模拟\n（单缓冲/双缓冲、用户态↔内核态传递）',
     '乒乓缓冲 A/B 切换 + Canvas 粒子动画 + 用户缓冲区面板实时累加'],
    ['④', '异常场景模拟\n（权限不足、非法地址、设备故障等）',
     '5种故障注入：EACCES / EFAULT / EIO / EPERM / ENOENT，运行时切换'],
    ['⑤', '单步执行，清晰展示每步细节',
     'btn-step 单步 + 自动连点 (50ms) + 系统调用日志 (13px) + 层级详情页'],
]
add_table(slide, 0.8, 2.0, col_w, headers, rows, Pt(12))

# ══════════════════════════════════════════
# P3 · 我们的目标
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '我们的目标', '不只是「交作业」—— 做一个真正能帮助理解 OS I/O 体系的教学工具')
add_page_number(slide, 3)

cards = [
    ('🖥 浏览器内可交互', '每一步都可暂停、观察、回溯\n单步与自动连点两种演示模式'),
    ('🔍 全链路可视化', '数据从用户空间到磁盘控制器再返回\nIRP↓ / IRQ↑ / DATA↑ 三条路径全程可观测'),
    ('⚡ 多种异常可切换注入', 'EACCES / EFAULT / EIO / EPERM / ENOENT\n抽象错误码变成可见的拦截过程'),
]
for i, (title, body) in enumerate(cards):
    x = 0.8 + i * 4.0
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(2.2), Inches(3.6), Inches(2.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(16)
    tf.margin_top = Pt(16)
    tf.margin_right = Pt(12)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PHOSPHOR
    p.font.name = '微软雅黑'
    p.space_after = Pt(16)

    for line in body.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = INK_DIM
        p.font.name = '宋体'
        p.space_after = Pt(6)

# ══════════════════════════════════════════
# P4 · 最简模型设计
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '最简模型设计', '课本四层 I/O 软件 + 一层硬件 = 五层全链路')
add_page_number(slide, 4)

layers = [
    ('L0', '用户层 I/O 软件', 'sys_read(fd, buf, len)', PHOSPHOR),
    ('L1', '设备无关软件 (VFS)', '路径解析 → 权限校验 → 缓冲分配', COBALT),
    ('L2', '设备驱动程序', '构造 IRP → 写控制寄存器 → 进程阻塞', AMBER),
    ('L4', '硬件层 (磁盘控制器)', '磁盘寻道 → 扇区读取 → DMA 传输', RGBColor(0x88, 0x88, 0x88)),
    ('L3', '中断处理程序 (ISR)', 'ISR 响应 → 数据拷贝 → 唤醒进程', GREEN),
]
for i, (tag, name, desc, clr) in enumerate(layers):
    y = 1.8 + i * 1.02
    # tag
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.8), Inches(y), Inches(0.7), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = clr
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER

    txt_box(slide, 1.7, y + 0.02, 3.5, 0.4, name, Pt(17), INK, bold=True)
    txt_box(slide, 5.3, y + 0.02, 4.0, 0.4, desc, Pt(13), INK_DIM)

    if i < 4:
        arrow = 'IRP ↓' if i < 3 else 'IRQ ↑'
        arrow_clr = PHOSPHOR if i < 3 else GREEN
        txt_box(slide, 3.0, y + 0.55, 1.2, 0.35, arrow, Pt(12), arrow_clr, bold=True)

# Right side: note
add_bullet_card(slide, 9.2, 1.8, 3.6, 2.2,
                '最简模型说明',
                ['此时没有双缓冲', '没有 ACL 权限校验', '没有故障注入', '仅正常流转'],
                COBALT)

add_bullet_card(slide, 9.2, 4.3, 3.6, 2.5,
                'IRP / IRQ 模型',
                ['IRP↓: I/O Request Packet 下发', 'IRQ↑: Interrupt Request 硬件中断',
                 '进程 BLOCKED 等待中断', 'ISR 唤醒进程 + copy_to_user'],
                PHOSPHOR)

# ══════════════════════════════════════════
# P5 · 前后端架构选型
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '前后端架构选型', '技术选型的「为什么」比「用什么」更重要')
add_page_number(slide, 5)

headers = ['决策点', '选择', '理由']
col_w = [2.0, 3.0, 6.5]
rows = [
    ['后端语言', 'Go', '单二进制部署、gRPC 原生支持、交叉编译方便（GOOS=linux CGO_ENABLED=0）'],
    ['通信协议', 'gRPC-Web 双向流', '前端点一步 → 后端推一个快照 → 天然匹配状态机模型\nWebSocket bidi streaming，无需轮询'],
    ['前端框架', 'Vanilla JS', '无 React/Vue 依赖，直接操作 DOM 更透明\n零 npm 依赖（除 google-protobuf 运行时）'],
    ['图形引擎', 'SVG + Canvas', 'SVG 做拓扑连线（CSS 动画联动），Canvas 做粒子系统（60fps 贝塞尔曲线）'],
    ['部署方式', '二进制 + systemd + NPM', '无容器依赖，单文件上传即运行\nNginx Proxy Manager 反代 + WebSocket Upgrade'],
]
add_table(slide, 0.8, 2.0, col_w, headers, rows, Pt(12))

# Architecture diagram
add_code_block(slide, 0.8, 4.8, 11.5, 2.2,
               """  浏览器 (SVG + Canvas + JS)  ←──WebSocket bidi──→  Go 后端
   ├─ 五层拓扑图 + 6条有向连线           ├─ gRPC-Web Server
   ├─ Canvas 粒子系统 (requestAnimationFrame)  ├─ 状态机引擎 (NextStep)
   ├─ 手风琴信息面板 (进程/内存/硬件)         │   ├─ 五层递进推进
   ├─ 层级详情页 (detail.html?layer=XX)      │   ├─ 双缓冲乒乓切换
   └─ gRPC-Web 帧编解码器 (手写120行)         │   └─ 5种故障注入
                                          └─ ACL 四级校验管线""")

# ══════════════════════════════════════════
# P6 · 最简原型
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '最简原型', '第一版：五层卡片 + 一条线 + 「下一步」按钮')
add_page_number(slide, 6)

add_bullet_card(slide, 0.8, 1.8, 5.8, 2.5,
                '第一版功能',
                ['五层卡片 + 1 条 SVG 连线', '「下一步」按钮推动一层',
                 '仅正常流转：read → VFS → 驱动 → 硬件 → ISR → 返回',
                 '通信：前端 WebSocket ↔ Go 后端 gRPC-Web'],
                PHOSPHOR)

add_bullet_card(slide, 7.2, 1.8, 5.5, 2.5,
                '这一版验证了什么',
                ['前后端通信链路通畅', 'gRPC-Web 双向流帧格式正确',
                 '状态机模型（5层递进）正确', 'Protobuf 序列化/反序列化正常',
                 'DOM 驱动渲染可行'],
                COBALT)

add_bullet_card(slide, 0.8, 4.8, 11.9, 2.2,
                '从零到一的关键决策',
                ['不引入任何前端框架 — 教学项目的代码应该让组员和老师一眼看懂',
                 '手写 gRPC-Web WebSocket 帧编解码器 — 120行代码替代 2万行 npm 依赖，对整个通信路径有完全控制',
                 'Go 后端单文件部署 — 交叉编译后直接 scp 上传运行，不需要 Go 工具链或容器环境'],
                INK_DIM)

# ══════════════════════════════════════════
# P7 · 为什么加入双缓冲
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '为什么加入双缓冲', '参照课本「通道」的设计思想 → 选择更适合教学场景的「双缓冲」方案')
add_page_number(slide, 7)

# Evolution line
evolutions = ['程序直接控制\n(轮询)', '中断驱动\nI/O', 'DMA\n控制器', 'I/O 通道\n处理器']
for i, (name, desc) in enumerate([(n, '') for n in evolutions]):
    x = 0.8 + i * 3.1
    clr = COBALT if i < 3 else (PHOSPHOR if i == 3 else GRAY_LIGHT)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(2.0), Inches(2.8), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = clr
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = INK
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER

    if i == 3:
        txt_box(slide, x + 0.3, 3.35, 2.2, 0.4, '← 我们最初考虑过', Pt(10), PHOSPHOR, bold=True)

    if i < 3:
        # arrow
        shape_a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                         Inches(x + 2.85), Inches(2.35), Inches(0.22), Inches(0.5))
        shape_a.fill.solid()
        shape_a.fill.fore_color.rgb = BORDER
        shape_a.line.fill.background()

# Comparison table
headers = ['', 'I/O 通道 (Channel)', '双缓冲 (Double Buffering)']
col_w = [1.8, 4.8, 4.8]
rows = [
    ['并行的层面', '通道程序独立于 CPU 执行\n通过 CCW 管理多台设备', '两块缓冲区交替读写\nDMA 写一块时 CPU 读另一块'],
    ['实现复杂度', '高：需要 CCW 状态机\nCAW/CSW 握手、多设备仲裁', '中：状态机直观建模\n前端可直接可视化切换过程'],
    ['教学价值', '偏硬件体系结构\n超出选题七范围', '偏 OS 软件层\n与课本缓冲技术章节直接对应'],
    ['可视化效果', '难以在浏览器呈现', '缓冲 A/B 颜色高亮 + 粒子动画\n乒乓切换一目了然'],
]
add_table(slide, 0.8, 4.0, col_w, headers, rows, Pt(11))

# Conclusion
txt_box(slide, 0.8, 6.7, 11.5, 0.5,
        '结论：参照通道「减少 CPU 干预、I/O 与计算并行」的设计思想 → 选择双缓冲，在实现复杂度和教学效果之间取得最佳平衡',
        Pt(14), PHOSPHOR, bold=True, alignment=PP_ALIGN.CENTER, font_name='微软雅黑')

# ══════════════════════════════════════════
# P8 · 双缓冲的实现
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '双缓冲的实现', '核心算法：乒乓切换 (Ping-Pong Buffer Switching)')
add_page_number(slide, 8)

add_code_block(slide, 0.8, 1.8, 6.0, 5.0,
               """初始化:
  ActiveWriteBuffer = 1   // DMA 目标
  ActiveReadBuffer  = 1   // CPU 拷贝源
  TotalChunks = ceil(BytesToRead / 4096)
  CurrentChunk = 0

每次 ISR 响应后:
  if CurrentChunk + 1 < TotalChunks:
      CurrentChunk++
      ActiveReadBuffer  = ActiveWriteBuffer
      ActiveWriteBuffer = 3 - ActiveWriteBuffer  // 1↔2
      编程下次 DMA → 硬件写新 WriteBuffer
      进程 BLOCKED → 回到 Driver 层循环
  else:
      进程 RUNNING, read() 返回""")

add_bullet_card(slide, 7.4, 1.8, 5.3, 2.2,
                '前端可视化',
                ['Canvas 粒子从硬件卡飞向目标缓冲区',
                 '二次贝塞尔曲线轨迹 (控制点偏移50px)',
                 'requestAnimationFrame 60fps 循环',
                 '5个粒子，淡入淡出透明度'],
                PHOSPHOR)

add_bullet_card(slide, 7.4, 4.3, 5.3, 2.5,
                '关键参数',
                ['chunk 大小: 4KB (模拟磁盘扇区/页)',
                 '缓冲区数: 2个内核缓冲区 (A/B)',
                 '乒乓触发: ISR 检测还有剩余 chunk',
                 '进程状态: BLOCKED (等待 DMA 中断)',
                 '粒子触发: ActiveWriteBuffer 变更时'],
                COBALT)

# ══════════════════════════════════════════
# P9 · 为什么加入 ACL 权限校验
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '为什么加入 ACL 权限校验', '课本 VFS 层的核心职责 —— 四级安全校验管线')
add_page_number(slide, 9)

# Pipeline
steps = [
    ('① 路径穿越检测', '解析 .. 段 → 规范化\n判定是否逃逸主目录', RED),
    ('② 文件存在检查', 'FileSystemDB 查表\n返回 ENOENT', AMBER),
    ('③ 敏感文件保护', 'IsSensitive 标记\n仅 root (uid=0) 可读', PHOSPHOR),
    ('④ Unix 权限位', 'owner/group/other\n9 位 rwx 权限模型', COBALT),
]
for i, (title, desc, clr) in enumerate(steps):
    x = 0.8 + i * 3.1
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(1.8), Inches(2.8), Inches(2.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = clr
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = clr
    p.font.name = '微软雅黑'
    p.space_after = Pt(10)

    for line in desc.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = INK_DIM
        p.font.name = '宋体'
        p.space_after = Pt(4)

    if i < 3:
        shape_a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                         Inches(x + 2.85), Inches(2.7), Inches(0.22), Inches(0.5))
        shape_a.fill.solid()
        shape_a.fill.fore_color.rgb = BORDER
        shape_a.line.fill.background()

add_bullet_card(slide, 0.8, 4.7, 5.8, 2.3,
                '多身份支持',
                ['Root 超级用户 (uid=0): 绕过所有检查',
                 '普通用户 (uid=1000, gid=1000): 受权限位约束',
                 '前端下拉框切换用户身份 → 同文件不同结果'],
                PHOSPHOR)

add_bullet_card(slide, 7.2, 4.7, 5.5, 2.3,
                '路径穿越防御',
                ['../../etc/shadow → 解析规范化 → /etc/shadow',
                 '判断是否逃逸 /home/user1/ 主目录',
                 '逃逸 → EPERM，不逃逸 → 继续校验'],
                RED)

# ══════════════════════════════════════════
# P10 · ACL 实现细节
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, 'ACL 实现细节', '模拟 Linux VFS 的文件系统元数据 + Unix 权限位检查')
add_page_number(slide, 10)

add_code_block(slide, 0.8, 1.8, 6.0, 2.8,
               """// FileSystemDB — 模拟文件系统
var FileSystemDB = map[string]*FileMetadata{
  "/etc/shadow":       {0, 0, 0600, true},   // root only
  "/home/user1/notes": {1000, 1000, 0644, false},
  "/tmp/data.bin":     {1000, 1000, 0666, false},
}

// Unix 权限位检查 (与课本一致)
func checkReadPermission(file, uid, gid):
   if uid == 0:       return true   // root 绕过
   if uid == owner:   check (perm>>6)&4  // owner read
   if gid == group:   check (perm>>3)&4  // group read
   return perm & 4                     // other read""")

add_code_block(slide, 7.4, 1.8, 5.3, 3.8,
               """// 四级校验管线
func ValidateFileAccess(userCtx, rawPath):
   // 1. 路径穿越检测
   if hasTraversalPattern(rawPath):
       resolved = resolvePath(rawPath)
       if escapesHomeDir(resolved, homeDir):
           return EFAULT  // 逃逸!

   // 2. 文件存在检查
   fileEntry = FileSystemDB[resolvePath(rawPath)]
   if !exists: return ENOENT

   // 3. 敏感文件检查
   if fileEntry.IsSensitive && uid != 0:
       return EACCES

   // 4. Unix 权限位检查
   if !checkReadPermission(fileEntry, uid, gid):
       return EACCES""")

add_bullet_card(slide, 0.8, 5.1, 11.9, 2.0,
                '校验失败处理',
                ['ValidationError{Code, Description} → 前端差错控制台实时展示',
                 'Code 映射标准 Unix 错误码: ENOENT / EACCES / EFAULT / EPERM',
                 'Description 为人类可读的描述，用于前端 stepDescription 和差错控制台'],
                PHOSPHOR)

# ══════════════════════════════════════════
# P11 · 故障注入体系
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '故障注入体系', '从「正常流转」到「什么情况都可能发生」—— 5种故障运行时切换')
add_page_number(slide, 11)

headers = ['故障类型', '错误码', '拦截层', '模拟场景']
col_w = [2.2, 1.8, 1.5, 6.0]
rows = [
    ['权限拒绝', 'EACCES', 'VFS 层', '普通用户 (uid=1000) 尝试读取 /etc/shadow 敏感系统文件'],
    ['非法地址', 'EFAULT', 'VFS 层', '用户传入 0xFFFFFFFF 内核空间地址，段错误越界'],
    ['硬件超时', 'EIO', '硬件层', '磁道损坏或设备掉线，STS_REG = DEVICE_ERROR'],
    ['路径穿越', 'EPERM', 'VFS 层', '使用 ../../etc/shadow 模式逃逸用户主目录范围'],
    ['文件不存在', 'ENOENT', 'VFS 层', '路径在 FileSystemDB 中不存在'],
]
add_table(slide, 0.8, 1.8, col_w, headers, rows, Pt(12))

add_bullet_card(slide, 0.8, 4.5, 5.8, 2.5,
                '运行时注入机制',
                ['前端下拉框选择故障类型', '发送 ACTION_INJECT_FAULT 指令',
                 '引擎在对应检查点拦截', '卡片变红 + 差错控制台展示'],
                PHOSPHOR)

add_bullet_card(slide, 7.2, 4.5, 5.5, 2.5,
                '教学价值',
                ['EACCES / EFAULT / EIO 不再是课本上抽象的字符串',
                 '亲眼看到进程在哪个阶段被拦截',
                 '理解错误码与 I/O 软件层的对应关系',
                 '可演示「同一路径 + 不同故障」的差异'],
                COBALT)

# ══════════════════════════════════════════
# P12 · 多身份检测演示
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '多身份检测演示', '同一文件，不同身份，不同结果')
add_page_number(slide, 12)

headers = ['测试场景', '用户', '文件', '预期结果', '原因']
col_w = [2.5, 2.0, 3.0, 1.8, 2.5]
rows = [
    ['Root 特权', 'root (uid=0)', '/etc/shadow', '✅ 通过', 'uid=0 绕过所有权限检查'],
    ['敏感文件拒绝', 'user1 (uid=1000)', '/etc/shadow', '❌ EACCES', '敏感文件仅 root 可读'],
    ['自己的文件', 'user1', '/home/user1/notes.txt', '✅ 通过', 'owner read bit = 4 (0644)'],
    ['公共文件', 'user1', '/tmp/data.bin', '✅ 通过', 'other read bit = 4 (0666)'],
]
add_table(slide, 0.8, 2.0, col_w, headers, rows, Pt(13))

add_bullet_card(slide, 0.8, 4.5, 11.9, 2.5,
                '演示要点',
                ['Root 用户 → 任何文件无障碍读取（超级用户特权，课本 §1.4.2）',
                 'user1 → 读取 /etc/shadow 被拒：先是敏感文件检查拦截（IsSensitive=true），即使绕过，0600 权限位也拒绝 other 读',
                 'user1 → 读取自己 home 目录下文件通过：owner match + owner read bit (0644 → 第6位=4)',
                 '/tmp/data.bin 权限 0666 → 所有用户可读写，与课本 Unix 权限模型一致'],
                PHOSPHOR)

# ══════════════════════════════════════════
# P13 · 最终系统架构
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '最终系统架构')
add_page_number(slide, 13)

add_code_block(slide, 0.8, 1.6, 11.5, 5.5,
               """  ┌── 浏览器前端 ─────────────────────────────┐      ┌── Go 后端 ────────────────────────────┐
  │                                                │      │                                        │
  │  SVG 五层拓扑图 (viewBox 720×750)               │      │  cmd/server/main.go                     │
  │  ├─ 6条有向连线: IRP↓ / IRQ↑ / DATA↑           │      │  ├─ grpcweb.WrapServer                   │
  │  ├─ 层卡片居中: left:50% + translateX(-50%)    │      │  └─ http.FileServer (STATIC_DIR)         │
  │  └─ foreignObject 标签联动激活                  │      │                                        │
  │                                                │      │  internal/service/handler.go             │
  │  Canvas 粒子系统                                │      │  ├─ ACTION_INIT → NewEngine + Send init  │
  │  ├─ getBoundingClientRect 动态坐标              │      │  ├─ ACTION_STEP_NEXT → NextStep + Send   │
  │  ├─ 二次贝塞尔曲线 (控制点偏移50px)              │      │  └─ ACTION_INJECT_FAULT → 运行时注入     │
  │  └─ requestAnimationFrame 60fps               │      │                                        │
  │                                                │      │  internal/engine/engine.go               │
  │  手风琴信息面板                                  │      │  ├─ 五层递进 NextStep 状态机             │
  │  ├─ 进程状态: RUNNING / BLOCKED / READY        │      │  ├─ 双缓冲乒乓切换 (ActiveWrite 1↔2)     │
  │  ├─ 用户缓冲区: 实时累加显示                     │      │  ├─ 5种故障注入拦截点                    │
  │  ├─ 差错控制台: 错误码 + 描述                    │      │  └─ 进程状态: RUNNING→BLOCKED→RUNNING    │
  │  └─ 系统调用日志: 80条循环缓冲                   │      │                                        │
  │                                                │      │  internal/engine/filesystem.go            │
  │  层级详情页 detail.html?layer=XX                │      │  ├─ FileSystemDB 文件元数据              │
  │  ├─ 6:4 Grid 布局                              │      │  └─ 四级校验管线 (穿越→存在→敏感→权限)   │
  │  ├─ 每层6步操作 + 伪代码                        │      │                                        │
  │  └─ 迷你拓扑位置图                               │      └────────────────────────────────────────┘
  │                                                │            ↑↓ WebSocket bidi streaming
  │  grpc-entry.js (手写帧编解码器)                  │            [control=0][flag][4B BE len][pb]
  │  [control=0][flag][4B BE len][pb data]         │
  └────────────────────────────────────────────────┘""")

# ══════════════════════════════════════════
# P14 · 核心功能一览
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '核心功能一览')
add_page_number(slide, 14)

features = [
    ('五层拓扑图 + 6条有向连线', 'IRP↓ 下发路径 / IRQ↑ 中断路径 / DATA↑ 数据返回路径\nSVG viewBox 坐标系 + 绝对定位卡片', PHOSPHOR),
    ('双缓冲乒乓切换', 'A/B 缓冲交替读写，DMA 写一块时 CPU 读另一块\nCanvas 粒子动画：贝塞尔曲线动态飞向目标', COBALT),
    ('5种故障注入', '运行时下拉切换，每种的拦截点不同\nEACCES / EFAULT / EIO / EPERM / ENOENT', RED),
    ('多身份 ACL', 'Root(uid=0) vs user1(uid=1000)\n四级校验管线：路径穿越→存在→敏感→权限位', AMBER),
    ('层级详情页', '点击卡片跳转 detail.html?layer=XX\n每层：输入输出 + 数据结构 + 6步操作 + 伪代码', COBALT),
    ('手风琴信息面板', '进程状态 / 用户缓冲区 / 差错控制台 / 系统调用日志\nflex 弹性链自动分配垂直空间', PHOSPHOR),
    ('单步 / 自动连点', '单步：暂停观察每层状态\n自动：50ms 间隔连续演示完整流转', INK),
]
for i, (title, desc, clr) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.8 + row * 2.7

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(3.8), Inches(2.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = clr
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = clr
    p.font.name = '微软雅黑'
    p.space_after = Pt(8)

    for line in desc.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = INK_DIM
        p.font.name = '宋体'
        p.space_after = Pt(3)

# ══════════════════════════════════════════
# P15 · 技术难点
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '技术难点')
add_page_number(slide, 15)

challenges = [
    ('gRPC-Web WebSocket 帧格式', '手写编解码器',
     ['Improbable 协议: [control=0][flag][4B BE len][pb data]',
      '120行代码替代 2万行 npm 依赖',
      '处理粘包缓冲 (parserBuf 累加 + 逐帧解析)',
      '消息队列: WebSocket 握手完成前缓存消息',
      'ASCII header 块作为握手首帧发送']),
    ('Canvas 粒子动画坐标', '动态坐标系统',
     ['SVG viewBox 虚拟坐标 vs Canvas 像素坐标',
      'getBoundingClientRect() 运行时动态计算起终点',
      '二次贝塞尔曲线: 控制点偏移50px 形成弧线',
      'requestAnimationFrame 60fps 循环',
      '淡入淡出透明度: t<0.15 淡入, t>0.85 淡出']),
    ('CSS 多层弹性布局', 'flex 弹性链',
     ['右侧4个手风琴面板: flex 纵向布局',
      '关键技巧: flex:1 + height:0 + min-height:0',
      '强制从零分配剩余空间，防止内容溢出',
      '手风琴折叠/展开不破坏 flex 链',
      '页尾 32px + 顶栏 72px 精确高度管控']),
    ('层级卡片居中', '任意宽度精确居中',
     ['left:50% ▸ 相对于包含块宽度',
      'translateX(-50%) ▸ 相对于元素自身宽度',
      '五种不同宽度卡片 (330px~560px) 均精确居中',
      '不受 position:absolute 和父元素影响']),
]
for i, (title, subtitle, bullets) in enumerate(challenges):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.6 + row * 2.9

    add_bullet_card(slide, x, y, 5.8, 2.6, f'{title}\n{subtitle}', bullets,
                    PHOSPHOR if col == 0 else COBALT)

# ══════════════════════════════════════════
# P16 · 测试覆盖
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '测试覆盖')
add_page_number(slide, 16)

add_bullet_card(slide, 0.8, 1.6, 5.8, 2.5,
                '正常流测试 (3场景)',
                ['单缓冲 + 普通文件 → 完整5步通过',
                 '双缓冲 + 普通文件 → 乒乓切换可见，粒子动画触发',
                 '双缓冲 + 64KB大文件 → 16个4KB chunk完整传输，A/B多次翻转'],
                PHOSPHOR)

add_bullet_card(slide, 7.2, 1.6, 5.5, 2.5,
                '故障注入测试 (5场景)',
                ['EACCES: user1 读 /etc/shadow → VFS拦截',
                 'EFAULT: 0xFFFFFFFF 地址 → VFS拦截',
                 'EIO: fault=3 + 正常路径 → 硬件层拦截',
                 'EPERM: ../../etc/shadow → VFS拦截',
                 'ENOENT: /nonexistent → VFS拦截'],
                RED)

add_bullet_card(slide, 0.8, 4.5, 5.8, 2.5,
                '权限边界测试 (4场景)',
                ['Root 读 /etc/shadow → ✅ 通过',
                 'user1 读 /etc/shadow → ❌ EACCES',
                 'user1 读自己的 notes.txt → ✅ 通过',
                 'user1 读 /tmp/data.bin → ✅ 通过'],
                COBALT)

add_bullet_card(slide, 7.2, 4.5, 5.5, 2.5,
                '前端集成测试 (6项)',
                ['连接状态机联动 (灯色 + 按钮启用)',
                 '单步 / 自动连点无卡顿',
                 '手风琴折叠不破坏布局',
                 '详情页跳转 + 返回按钮',
                 'DATA↑ 虚线有数据时点亮',
                 '粒子动画缩放后无偏移'],
                AMBER)

# ══════════════════════════════════════════
# P17 · 小组分工
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '小组分工 (8人)')
add_page_number(slide, 17)

headers = ['角色', '人数', '负责内容']
col_w = [3.0, 0.8, 7.8]
rows = [
    ['项目架构 + 核心开发 (组长)', '1', '整体架构设计、技术选型、Go 后端全部代码、前端全部代码、gRPC-Web 帧编解码器、云端部署、README + DESIGN + SLIDES 文档'],
    ['状态机引擎', '1', '五层递进推进 (NextStep)、双缓冲乒乓切换控制器、chunk 分块 (4KB) DMA 调度、进程状态转换 (RUNNING→BLOCKED→RUNNING)、5种故障注入拦截点'],
    ['文件系统校验', '1', 'ACL 四级管线 (路径穿越→存在→敏感→权限位)、FileSystemDB 元数据设计、Unix 9位权限模型实现、ValidationError 错误体系'],
    ['Proto 协议与通信', '1', 'Protobuf 8 message + 1 service 定义、gRPC-Web 双向流对接、WebSocket 帧格式编解码、消息队列 + 粘包缓冲处理'],
    ['前端可视化', '2', 'SVG 五层拓扑图 + 6条有向连线动画 (IRP↓/IRQ↑/DATA↑)、Canvas 粒子系统 (贝塞尔曲线 + 60fps)、手风琴信息面板、层级详情页内容填充、快照渲染与交互控制'],
    ['前端样式与布局', '1', '6:4 CSS Grid 宽屏排版、CSS 变量体系、浅色实验室主题、卡片居中定位体系、多页面导航 (detail.html)、响应式适配'],
    ['部署与文档', '1', '跨平台编译 + systemd + NPM 云端部署、测试用例设计与执行、README/DESIGN/SLIDES 文档协作、答辩 PPT 制作与现场 Demo'],
]
add_table(slide, 0.5, 1.7, col_w, headers, rows, Pt(12))

# ══════════════════════════════════════════
# P18 · 收获与展望
# ══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, '收获与展望')
add_page_number(slide, 18)

add_bullet_card(slide, 0.8, 1.6, 5.8, 2.8,
                '收获',
                ['对 OS I/O 软件层的理解从课本文字变成了\n可运行、可逐步观察的系统',
                 '掌握了 gRPC 双向流在实际项目中的应用',
                 '通过手写帧编解码器理解了 RPC 协议栈底层',
                 'CSS Grid/Flexbox 复杂布局实战经验',
                 'Canvas 动画系统 + 状态同步的开发经验',
                 '全栈独立开发能力 (Go + JS + 部署)'],
                PHOSPHOR)

add_bullet_card(slide, 7.2, 1.6, 5.5, 2.8,
                '可扩展方向',
                ['写操作模拟: sys_write 全链路',
                 'I/O 调度算法可视化: 电梯算法 / FCFS',
                 '多进程并发 I/O 竞争',
                 '磁盘调度动画: SCAN / C-SCAN',
                 '更多 I/O 控制方式: 通道模拟',
                 '性能分析面板: I/O 吞吐 / 延迟'],
                COBALT)

# Thank you
txt_box(slide, 0.8, 5.2, 11.5, 0.8, '谢谢！', Pt(36), INK, bold=True, alignment=PP_ALIGN.CENTER, font_name='微软雅黑')
txt_box(slide, 0.8, 5.9, 11.5, 0.5, 'https://github.com/smithluo11/os-read-', Pt(13), COBALT, alignment=PP_ALIGN.CENTER)
txt_box(slide, 0.8, 6.5, 11.5, 0.5, '操作系统课程设计 · 选题七 · I/O 软件层 read 操作模拟', Pt(12), INK_DIM, alignment=PP_ALIGN.CENTER)

# ── Save ──
output_path = '答辩PPT_OS_IO_Simulator.pptx'
prs.save(output_path)
print(f'Done! Saved to {output_path}')
